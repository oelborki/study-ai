import sys
sys.stdout.reconfigure(encoding="utf-8")
import json
from pptx import Presentation

def extract(path: str):
    prs = Presentation(path)
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    texts.append(t)

        title = texts[0] if texts else ""
        bullets = texts[1:] if len(texts) > 1 else []

        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""

        slides.append({
            "index": i,
            "title": title,
            "bullets": bullets,
            "notes": notes
        })

    return {"slides": slides}

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "missing pptx path"}))
        sys.exit(1)

    path = sys.argv[1]
    print(json.dumps(extract(path), ensure_ascii=False))
